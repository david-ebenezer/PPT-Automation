from pptx import Presentation

def add_verse(reference,verse):
    #Verse

    title_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = verse #verse
    subtitle.text = reference # reference 

def add_subpoint(Tamil,English):
    #SubSubpoints 

    title_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = English # English 

prs = Presentation('newtemplate.pptx')
docfile = open("myfile.txt",encoding="utf8")
doclines=docfile.readlines()
questions="".join(doclines).split("\n\n")

for i in questions:
    qindex=i.find("\n")
    add_subpoint("",i[:qindex])
    add_subpoint("",i[qindex:])
    
prs.save('testnew.pptx')
print("Done")

'''

add_subpoint("தடைகளை நீக்கிப்போடுகிறவர் ", 'The one who breaks ')
if ":" in f[j]:
v=True

'''
